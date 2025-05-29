import streamlit as st
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


from io import BytesIO

# Initialize Groq client
def initialize_groq_client():
    return Groq(api_key="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")

# Theme templates
THEMES = {
    "Royal Blue": {
        "background": RGBColor(0, 32, 96),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(255, 255, 255),
        "font": "Calibri"
    },
    "Modern Green": {
        "background": RGBColor(214, 234, 223),
        "title_color": RGBColor(0, 102, 51),
        "text_color": RGBColor(51, 51, 51),
        "font": "Arial"
    },
    "Creative Orange": {
        "background": RGBColor(255, 243, 224),
        "title_color": RGBColor(230, 115, 0),
        "text_color": RGBColor(64, 64, 64),
        "font": "Verdana"
    },
    "Elegant Purple": {
        "background": RGBColor(242, 240, 255),
        "title_color": RGBColor(102, 51, 153),
        "text_color": RGBColor(51, 51, 51),
        "font": "Georgia"
    }
}

def rgbcolor_to_tuple(rgb_color):
    """Convert RGBColor object to RGB tuple"""
    if hasattr(rgb_color, 'rgb'):  # For RGBColor objects
        # RGBColor stores color as a long integer in format 0x00RRGGBB
        rgb_int = rgb_color.rgb
        r = (rgb_int >> 16) & 0xFF
        g = (rgb_int >> 8) & 0xFF
        b = rgb_int & 0xFF
        return r, g, b
    elif isinstance(rgb_color, (tuple, list)) and len(rgb_color) == 3:  # Already a tuple
        return rgb_color
    else:  # Fallback
        return (0, 0, 0)  # Default black if conversion fails

def show_theme_previews():
    """Display theme selection cards with proper color conversion"""
    st.write("")  # Add spacing
    
    cols = st.columns(len(THEMES))
    selected_theme = st.session_state.get("selected_theme", "Royal Blue")
    
    for i, (theme_name, theme) in enumerate(THEMES.items()):
        with cols[i]:
            # Convert RGBColor to RGB tuples
            bg_r, bg_g, bg_b = rgbcolor_to_tuple(theme['background'])
            title_r, title_g, title_b = rgbcolor_to_tuple(theme['title_color'])
            text_r, text_g, text_b = rgbcolor_to_tuple(theme['text_color'])
            
            # Highlight selected theme
            border = "3px solid #555" if theme_name == selected_theme else "none"
            
            st.markdown(
                f"""
                <div style="background-color:rgb({bg_r},{bg_g},{bg_b});
                    padding:15px;border-radius:8px;height:110px;
                    border:{border};margin-bottom:10px;
                    cursor:pointer;" onclick="document.getElementById('theme_{i}').click()">
                    <h4 style="color:rgb({title_r},{title_g},{title_b});
                        margin-top:0;margin-bottom:8px;font-size:16px;">
                    {theme_name}</h4>
                    <p style="color:rgb({text_r},{text_g},{text_b});
                        margin-bottom:0;font-size:14px;">
                    Sample content</p>
                </div>
                """,
                unsafe_allow_html=True
            )
            
            # Hidden radio button for selection
            st.radio(
                "Theme",
                options=[theme_name],
                key=f"theme_{i}",
                index=0 if theme_name == selected_theme else None,
                label_visibility="collapsed",
                on_change=lambda: st.session_state.update(selected_theme=st.session_state[f"theme_{i}"])
            )

# Generate slide content with Groq

def generate_slide_content(topic, slide_count):
    """Generate presentation content using specified Groq model"""
    client = Groq(api_key="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")
    
    response = client.chat.completions.create(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        messages=[
            {
                "role": "system",
                "content": """You MUST create PowerPoint content in this EXACT format:
                
                ---SLIDE 1---
                TITLE: [Presentation Title]
                CONTENT: [Content for this slide]
                
                ---SLIDE 2---
                TITLE: [Slide Title] 
                CONTENT: [Content for this slide]
                
                [Continue for all slides]
                
                Rules:
                1. Use ---SLIDE X--- dividers
                2. Always include both TITLE: and CONTENT: labels
                3. Generate 4 points on each slide and use • as bullet points
                4. Keep content concise but informative"""
            },
            {
                "role": "user",
                "content": f"Create a {slide_count}-slide PowerPoint about {topic} with these slides:\n"
                           "1. Title Slide (include topic)\n"
                           "2. Introduction\n"
                           "3. Key Feature 1\n" 
                           "4. Key Feature 2\n"
                           "5. Conclusion\n"
                           "Use the exact format specified above."
            }
        ],
        temperature=0.7,
        max_tokens=4000,
        stream=False
    )
    
    return response.choices[0].message.content

def parse_slides(input_text):
    slides = []
    chunks = [chunk.strip() for chunk in input_text.split('---') if chunk.strip()]
    
    for chunk in chunks:
        lines = chunk.splitlines()
        title = next((line.split(":", 1)[1].strip() for line in lines if line.strip().startswith("TITLE:")), "")
        
        # Find index of "CONTENT:" line
        content_index = next((i for i, line in enumerate(lines) if line.strip().startswith("CONTENT:")), None)
        content = ""
        if content_index is not None:
            content_lines = lines[content_index + 1:]
            content = " ".join(line.strip() for line in content_lines if line.strip())

        if title and content:
            slides.append({"title": title, "content": content})

    # Output the result
    for slide in slides:
        print(f'    {{"title": "{slide["title"]}", "content": "{slide["content"]}"}}')
    
    return slides



def create_presentation(slides_data, theme_name):
    theme = THEMES[theme_name]
    prs = Presentation()
    
    # Set slide size (16:9)
    prs.slide_width = Inches(12)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    for i, slide in enumerate(slides_data):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide_layout = prs.slides.add_slide(layout)
        
        # Background
        background = slide_layout.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme["background"]
        
        # Title
        title = slide_layout.shapes.title
        title.text = slide["title"]
        title.text_frame.paragraphs[0].font.color.rgb = theme["title_color"]
        title.text_frame.paragraphs[0].font.name = theme["font"]
        title.text_frame.paragraphs[0].font.size = Pt(36 if i == 0 else 32)
        title.text_frame.paragraphs[0].font.bold = True

# Move title 0.5 inches to the right (adjust the value as needed)
        # Content
        
        if i > 0:
            content = slide_layout.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()  # Clear existing content

        # Split the content into points using '•' as delimiter
            bullet_points = [point.strip() for point in slide["content"].split("•") if point.strip()]

            for idx, point in enumerate(bullet_points):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = f"{point}"
                p.level = 0  # Bullet level (0 for top-level)
                p.font.color.rgb = theme["text_color"]
                p.font.name = theme["font"]
                p.font.size = Pt(18)
        
    # Save to bytes
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def main():
    st.set_page_config(page_title="AI PowerPoint Generator", page_icon="📊", layout="centered")
    
    st.title("🎨 AI PowerPoint Generator")
    st.markdown("Create professional presentations in seconds")
    st.write("")  # Add spacing
    
    # User inputs
    with st.container():
        col1, col2 = st.columns(2)
        with col1:
            topic = st.text_input("**Presentation Topic**", "The Future of AI")
        with col2:
            slide_count = st.slider("**Number of Slides**", 3, 10, 5)
    
    # Theme selection
    st.subheader("🎨 Select a Theme")
    show_theme_previews()
    
    # Get selected theme
    selected_theme = next(
        (name for i, name in enumerate(THEMES) 
         if st.session_state.get(f"theme_{i}") == name),
        list(THEMES.keys())[0]
    )
    
    # Generate button
    if st.button("✨ Generate Presentation", type="primary", use_container_width=True):
        with st.spinner("Creating your presentation..."):
            try:
                ai_content = generate_slide_content(topic, slide_count)
                slides_data = parse_slides(ai_content)
                
                pptx = create_presentation(slides_data, selected_theme)
                
                st.success("✅ Presentation generated successfully!")
                
                # Download and preview
                dl_col, prev_col = st.columns([1, 2])
                
                with dl_col:
                    st.download_button(
                        label="📥 Download PowerPoint",
                        data=pptx,
                        file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                
                with st.expander("Preview Content"):
                    for i, slide in enumerate(slides_data):
                        st.subheader(f"Slide {i+1}: {slide['title']}")
                        st.write(slide["content"])
                        
                            
            except Exception as e:
                st.error(f"❌ Error: {str(e)}")

if __name__ == "__main__":
    main()